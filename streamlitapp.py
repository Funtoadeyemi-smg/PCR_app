from pptx import Presentation
import re
import os
import json
from datetime import date
import pandas as pd
from time import sleep
from openai import OpenAI
from dotenv import load_dotenv
import streamlit as st
import tempfile
import os
from PIL import Image
import base64
import matplotlib.pyplot as plt
from datetime import datetime
from pptx.util import Inches
from pptx.oxml import parse_xml
from pptx.oxml.ns import qn

load_dotenv()
api_key = os.getenv('OPENAI_API_KEY')
connection_string = os.getenv('CONNECTION_STRING')

client = OpenAI(api_key=api_key)

class PowerPointProcessor:
    """Handles PowerPoint modifications."""
    
    def __init__(self, pptx_file):
        self.pptx_file = pptx_file
        self.prs = Presentation(pptx_file)

    def replace_placeholders(self, replacements, is_pin, output_pptx):
        """
        Replaces placeholders in a PowerPoint presentation with provided values, including text inside tables.
        
        Parameters:
        - replacements (dict): Dictionary mapping placeholders to actual values.
        - output_pptx (str): Path where the updated PowerPoint should be saved.
        """
        found_placeholders = set()  # Store found placeholders for debugging

        for slide in self.prs.slides:
            # Replace in normal text boxes
            for shape in slide.shapes:
                if shape.has_text_frame:
                    for paragraph in shape.text_frame.paragraphs:
                        for run in paragraph.runs:
                            for placeholder, replacement in replacements.items():
                                if placeholder in run.text:
                                    found_placeholders.add(placeholder)
                                    run.text = run.text.replace(placeholder, str(replacement))

            for shape in slide.shapes:
                if shape.has_table:
                    table = shape.table
                    for row in table.rows:
                        for cell in row.cells:
                            for paragraph in cell.text_frame.paragraphs:
                                for run in paragraph.runs:
                                    for placeholder, replacement in replacements.items():
                                        if placeholder in run.text:
                                            found_placeholders.add(placeholder)
                                            run.text = run.text.replace(placeholder, str(replacement))

            for shape in slide.shapes:
                if shape.has_table:
                    table = shape.table
                    rows_to_remove = []
                    for i, row in enumerate(table.rows):
                        texts = [cell.text.strip() for cell in row.cells]
                        if all("{" in text and "}" in text for text in texts):  # all look like placeholders
                            rows_to_remove.append(i)

                    for i in reversed(rows_to_remove):  # delete from bottom up to keep index stable
                        tbl = table._tbl
                        tr = tbl.tr_lst[i]
                        tbl.remove(tr)

        try:
            chart_slide = self.prs.slides[8]  # Slide 10
            for shape in chart_slide.shapes:
                if shape.shape_type == 13:  # MSO_SHAPE_TYPE.PICTURE
                    chart_slide.shapes._spTree.remove(shape._element)
            chart_slide.shapes.add_picture("impressions_chart.png", Inches(0.7), Inches(2.9), height=Inches(3.2))
            chart_slide.shapes.add_picture("clicks_chart.png", Inches(8.7), Inches(2.9), height=Inches(3.2))

        except IndexError:
            print("⚠️ Slide 10 (index 9) not found – skipping chart replacement.")
        except FileNotFoundError as e:
            print(f"⚠️ Image file not found: {e.filename}")    

        if not is_pin:
            # Delete slides 27, 28, and 29 by index (in reverse order to avoid reindexing issues)
            for i in [28, 27, 26, 25]:
                try:
                    xml_slides = self.prs.slides._sldIdLst
                    slide_id = xml_slides[i]
                    xml_slides.remove(slide_id)
                except IndexError:
                    print(f"⚠️ Slide index {i+1} not found – skipping deletion.")


        self.prs.save(output_pptx)

        # Debugging output
        print(f"Updated presentation saved as: {output_pptx}")
        print("Placeholders found and replaced")

        if not found_placeholders:
            print("⚠️ No placeholders were replaced. Check if placeholders in the PPT match replacements keys.")


class DataExtractor:
    def __init__(self, meta_file, pinterest_file=None, media_plan_file=None, prompt_file=None):
        self.meta_file = meta_file
        self.pinterest_file = pinterest_file
        self.media_plan_file = media_plan_file
        self.prompt_file = prompt_file
        self.replacements = {}
        self.is_pin =True

    def open_file(self, filepath):
        with open(filepath, 'r', encoding='utf-8') as infile:
            return infile.read()
        
    def gpt4_completion(self, prompt, engine = 'gpt-4.1', temp=0, tokens=4000, response_format={ "type": "json_object" }):
        max_retry = 5
        retry = 0
        while True:
            try:
                response = client.chat.completions.create(
                    model=engine, 
                    messages = [
                {"role": "system", "content": "You are a helpful assistant designed to output JSON."},
                {"role": "user", "content": prompt}
            ],
                    temperature=temp,
                    max_tokens=tokens,
                    response_format=response_format,)
                
                text = response.choices[0].message.content.replace("\n", " ")
                return text
            except Exception as oops:
                retry += 1
                if retry >= max_retry:
                    return "GPT4 error: %s" % oops
                print('Error communicating with OpenAI:', oops)
                sleep(1)

    def plot_figures(self, est, actual, file_name):
        categories = ['Estimated', 'Actual']
        values = [est, actual]
        colors = ['#6FC0C3', '#4B88F1']

        # Create bar chart
        fig, ax = plt.subplots(figsize=(3.0, 2.5))
        bars = ax.bar(categories, values, color=colors)

        # Add white bold labels on bars
        for bar in bars: 
            yval = bar.get_height()
            ax.text(bar.get_x() + bar.get_width() / 2, yval - (yval * 0.05), f'{yval:,}', ha='center',va='top', color='white', fontweight='bold', fontsize=7)

        # Format y-axis with commas
        ax.get_yaxis().set_major_formatter(plt.FuncFormatter(lambda x, _: f'{int(x):,}'))
        # Add horizontal gridlines
        ax.yaxis.grid(True, linestyle='--', linewidth=0.5)
        ax.set_axisbelow(True)
        ax.spines['top'].set_visible(False)
        ax.spines['right'].set_visible(False)
        plt.xticks(fontsize=7)
        plt.yticks(fontsize=7)
        plt.tight_layout()
        plt.savefig(file_name, dpi=300)

    def extract_values(self):
        meta = pd.read_excel(self.meta_file)
        pinterest = pd.read_csv(self.pinterest_file) if self.pinterest_file else None
        media_plan = pd.read_excel(self.media_plan_file, skiprows=7, skipfooter=3) if self.media_plan_file else None
        prompt = self.open_file(self.prompt_file) if self.prompt_file else ""

        # ========== Meta data =============
        Brand_Name = meta['Brand'][1]
        Date_ = date.today().strftime('%d/%m/%Y')
        Campaign_Name = meta['A360_Campaign'][1]

        if pinterest is not None and meta is not None:
            channel_list_str = "Meta and Pinterest"
        elif meta is not None:
            channel_list_str = "Meta only"
        elif pinterest is not None:
            channel_list_str = "Pinterest only"
        else:
            channel_list_str = "No data available"

        gross_spend_meta = float(meta['Gross_Spend'].sum())
        net_spend_meta = gross_spend_meta * 0.7
        impressions_meta = int(meta['Impressions'].sum())
        clicks_meta = int(meta['Clicks'].sum())
        CTR_meta = (clicks_meta / impressions_meta) * 100
        net_CPM_meta = (net_spend_meta / impressions_meta) * 1000
        brand_revenue_meta = float(meta['Brand_Revenue'].sum())
        brand_ROAS_meta = brand_revenue_meta / net_spend_meta
        brand_ROI_meta = brand_revenue_meta / gross_spend_meta
        total_FSKU_rev = float(meta['SKU_Revenue'].sum())
        total_FSKU_sales = int(meta['SKU_Units'].sum())  
        total_FSKU_ROAS = total_FSKU_rev / (gross_spend_meta * 0.7)
        total_FSKU_ROI = total_FSKU_rev / gross_spend_meta
        brand_online_revenue_meta = float(meta['Brand_Online_Revenue'].sum())
        brand_instore_revenue_meta = float(meta['Brand_Instore_Revenue'].sum())
        FSKU_online_revenue_meta = float(meta['SKU_Online_Revenue'].sum())
        FSKU_instore_revenue_meta = float(meta['SKU_Instore_Revenue'].sum())
        gross_spend_per_audience_meta = meta.groupby('Ad_Set_Name')[['Gross_Spend']].sum().reset_index()
        net_spend_per_audience_meta = gross_spend_per_audience_meta['Gross_Spend'] * 0.7
        impressions_per_audience_meta = meta.groupby('Ad_Set_Name')[['Impressions']].sum().reset_index()
        reach_per_audience_meta = meta.groupby('Ad_Set_Name')[['Ad_Set_Reach']].sum().reset_index()
        clicks_per_audience_meta = meta.groupby('Ad_Set_Name')[['Clicks']].sum().reset_index()
        CTR_per_audience_meta = clicks_per_audience_meta['Clicks'] / impressions_per_audience_meta['Impressions'] * 100
        net_CPM_per_audience_meta = net_spend_per_audience_meta / impressions_per_audience_meta['Impressions'] * 1000
        brand_revenue_per_audience_meta = meta.groupby('Ad_Set_Name')[['Brand_Revenue']].sum().reset_index()
        gross_spend_series = gross_spend_per_audience_meta['Gross_Spend']
        brand_ROI_per_audience_meta = brand_revenue_per_audience_meta['Brand_Revenue'] / gross_spend_per_audience_meta['Gross_Spend']
        meta_input = meta.to_json(orient='records', indent=4)

        combined_meta_df = pd.DataFrame({
        'Ad_Set_Name': gross_spend_per_audience_meta['Ad_Set_Name'],
        'Gross_Spend': gross_spend_series,
        'Net_Spend': net_spend_per_audience_meta,
        'Impressions': impressions_per_audience_meta['Impressions'],
        'Reach': reach_per_audience_meta['Ad_Set_Reach'],
        'Clicks': clicks_per_audience_meta['Clicks'],
        'CTR': CTR_per_audience_meta,
        'Net_CPM': net_CPM_per_audience_meta,
        'Brand_Revenue': brand_revenue_per_audience_meta['Brand_Revenue'],
        'Brand_ROI': brand_ROI_per_audience_meta
        })

        #==============Pinterest data============
        if pinterest is not None:
            net_spend_pin = float(pinterest['Spend in account currency'].sum())
            gross_spend_pin = net_spend_pin / 0.7
            impressions_pin = int(pinterest['Impressions'].sum())
            clicks_pin = int(pinterest['Pin clicks'].sum())
            CTR_pin = (clicks_pin / impressions_pin) * 100
            net_CPM_pin = (net_spend_pin / impressions_pin) * 1000
            brand_revenue_pin = float(pinterest['Total order value (Lead)'].sum())
            brand_ROAS_pin = brand_revenue_pin / net_spend_pin
            brand_ROI_pin = brand_revenue_pin / gross_spend_pin
            brand_online_revenue_pin = float(pinterest['Web order value (Lead)'].sum())
            brand_instore_revenue_pin = float(pinterest['Offline order value (Lead)'].sum())
            net_spend_per_audience_pin = pinterest.groupby('Ad group name')[['Spend in account currency']].sum().reset_index()
            gross_spend_per_audience_pin = net_spend_per_audience_pin['Spend in account currency'] / 0.7
            impressions_per_audience_pin = pinterest.groupby('Ad group name')[['Impressions']].sum().reset_index()
            reach_per_audience_pin = pinterest.groupby('Ad group name')[['Reach']].sum().reset_index()
            clicks_per_audience_pin = pinterest.groupby('Ad group name')[['Pin clicks']].sum().reset_index()
            CTR_per_audience_pin = clicks_per_audience_pin['Pin clicks'] / impressions_per_audience_pin['Impressions'] * 100
            net_CPM_per_audience_pin = net_spend_per_audience_pin['Spend in account currency'] / impressions_per_audience_pin['Impressions'] * 1000
            brand_revenue_per_audience_pin = pinterest.groupby('Ad group name')[['Total order value (Lead)']].sum().reset_index()
            brand_ROI_per_audience_pin = brand_revenue_per_audience_pin['Total order value (Lead)'] / gross_spend_per_audience_pin
            gross_spend_series_pin = gross_spend_per_audience_pin
            pin_input = pinterest.to_json(orient='records', indent=4)

 
        else:
            self.is_pin = False
            net_spend_pin = gross_spend_pin = impressions_pin = clicks_pin = CTR_pin  = 0
            net_CPM_pin = brand_revenue_pin = brand_ROAS_pin =  brand_instore_revenue_pin =  brand_online_revenue_pin = brand_ROI_pin = 0
            pin_input = 'None'
            net_spend_per_audience_pin = impressions_per_audience_pin = reach_per_audience_pin = clicks_per_audience_pin= brand_revenue_per_audience_pin =pd.DataFrame(columns=['Ad group name', 'Spend in account currency', 'Impressions', 'Reach', 'Pin clicks', 'Total order value (Lead)'])
            gross_spend_per_audience_pin = gross_spend_series_pin = CTR_per_audience_pin = net_CPM_per_audience_pin = brand_ROI_per_audience_pin = pd.Series(dtype='float64')

        combined_pin_df = pd.DataFrame({
        'Ad_Set_Name': net_spend_per_audience_pin['Ad group name'],
        'Gross_Spend': gross_spend_series_pin,
        'Net_Spend': net_spend_per_audience_pin['Spend in account currency'],
        'Impressions': impressions_per_audience_pin['Impressions'],
        'Reach': reach_per_audience_pin['Reach'],
        'Clicks': clicks_per_audience_pin['Pin clicks'],
        'CTR': CTR_per_audience_pin,
        'Net_CPM': net_CPM_per_audience_pin,
        'Brand_Revenue': brand_revenue_per_audience_pin['Total order value (Lead)'],
        'Brand_ROI': brand_ROI_per_audience_pin
        })

        data3 = {
            "Metric": ["Brand_Revenue", "brand_ROAS", "brand_ROI"],
            "Meta": [brand_revenue_meta, brand_ROAS_meta, brand_ROI_meta],
            "Pin": [brand_revenue_pin, brand_ROAS_pin, brand_ROI_pin]
        }

        #==============total data============

        total_gross_spend = gross_spend_meta + gross_spend_pin
        total_impressions = impressions_meta + impressions_pin
        total_clicks = clicks_meta + clicks_pin
        total_CTR = total_clicks / total_impressions * 100 if total_impressions else 0
        total_CPM = (net_CPM_meta + net_CPM_pin) / 2 if pinterest is not None else net_CPM_meta
        data3 = {
            "Metric": ["Brand_Revenue", "brand_ROAS", "brand_ROI"],
            "Meta": [brand_revenue_meta, brand_ROAS_meta, brand_ROI_meta],
            "Pin": [brand_revenue_pin, brand_ROAS_pin, brand_ROI_pin]
        }
        channel_kpi = pd.DataFrame(data3)
        channel_kpi_comparison_json = channel_kpi.to_json(orient='records', indent=4)

        total_brand_revenue = brand_revenue_meta + brand_revenue_pin
        total_ROI = total_brand_revenue / total_gross_spend
        total_ROAS = total_brand_revenue / (total_gross_spend * 0.7)
        total_brand_online_rev = brand_online_revenue_meta + brand_online_revenue_pin
        total_brand_instore_rev = brand_instore_revenue_meta + brand_instore_revenue_pin

        data2 = {"total_brand_online_revenue": total_brand_online_rev, "total_brand_instore_revenue": total_brand_instore_rev}
        online_offline_json = json.dumps(data2, indent=4)
        merged_audience_metrics = pd.concat([combined_meta_df, combined_pin_df]).groupby('Ad_Set_Name', as_index=False).sum()
        cols_to_check = merged_audience_metrics.columns[1:]  # Skip first column
        merged_audience_metrics = merged_audience_metrics[((merged_audience_metrics[cols_to_check] == 0) | merged_audience_metrics[cols_to_check].isna()).sum(axis=1) / len(cols_to_check) < 0.8]
        merged_audience_metrics.drop(columns=['Brand_ROI', 'CTR', 'Net_CPM'], inplace=True)
        merged_audience_metrics['Brand_ROI'] = merged_audience_metrics['Brand_Revenue'] / merged_audience_metrics['Gross_Spend']
        merged_audience_metrics['CTR'] = merged_audience_metrics['Clicks'] / merged_audience_metrics['Impressions'] * 100
        merged_audience_metrics['Net_CPM'] = merged_audience_metrics['Net_Spend'] / merged_audience_metrics['Impressions'] * 1000
        audience_json = merged_audience_metrics.to_json(orient='records', indent=4)

        # ========== Media Plan Estimates =============
        if media_plan is not None:
            def parse_date_range(date_text):
                """
                Parse various date range formats including:
                - DD/MM/YYYY - DD/MM/YYYY
                - DD Month - DD Month YYYY
                - DDth Month - DDth Month YYYY
                """
                # Try standard format first: DD/MM/YYYY - DD/MM/YYYY
                standard_ranges = re.findall(r'(\d{1,2}/\d{1,2}/\d{4})\s*-\s*(\d{1,2}/\d{1,2}/\d{4})', date_text)
                if standard_ranges:
                    return [(datetime.strptime(start, "%d/%m/%Y"), datetime.strptime(end, "%d/%m/%Y")) 
                            for start, end in standard_ranges]
                
                # Try format: DDth Month - DDth Month YYYY
                # This pattern captures dates like "10th December - 23rd December 2024"
                month_names = "January|February|March|April|May|June|July|August|September|October|November|December"
                pattern = rf'(\d{{1,2}}(?:st|nd|rd|th)?\s+(?:{month_names}))\s*-\s*(\d{{1,2}}(?:st|nd|rd|th)?\s+(?:{month_names})(?:\s+\d{{4}})?)'
                
                textual_ranges = re.findall(pattern, date_text, re.IGNORECASE)
                
                if not textual_ranges:
                    return []
                
                parsed_dates = []
                for start_date, end_date in textual_ranges:
                    # Check if year is in end_date, if not try to extract from the original text
                    if not re.search(r'\d{4}', end_date):
                        # Try to find a year in the original text
                        year_match = re.search(r'\b(\d{4})\b', date_text)
                        year = year_match.group(1) if year_match else str(datetime.now().year)
                        # Add year to both start and end dates
                        start_date = f"{start_date} {year}"
                        end_date = f"{end_date}"  # Year should already be included
                    
                    # Clean up ordinal indicators (st, nd, rd, th)
                    start_date = re.sub(r'(\d+)(st|nd|rd|th)', r'\1', start_date)
                    end_date = re.sub(r'(\d+)(st|nd|rd|th)', r'\1', end_date)
                    
                    try:
                        # Parse the dates
                        start_dt = datetime.strptime(start_date, "%d %B %Y")
                        end_dt = datetime.strptime(end_date, "%d %B %Y")
                        parsed_dates.append((start_dt, end_dt))
                    except ValueError as e:
                        print(f"Error parsing date: {e}, start_date: {start_date}, end_date: {end_date}")
                        continue
                
                return parsed_dates

            # Usage in your existing code
            media_plan = media_plan[['Platform', 'Estimated Impressions', 'Estimated link clicks', 'Estimated Frequency', 'Estimated Reach', 'Gross spend by channel / platform', 'Estimated CTR', 'Net CPM', 'Flight duration']]
            media_plan['Platform'] = media_plan['Platform'].fillna(method='ffill')
            media_plan['Flight duration'] = media_plan['Flight duration'].fillna(method='ffill')
            #media_plan_clean = media_plan.dropna()
            media_plan_clean = media_plan.dropna(thresh=media_plan.shape[1] - 4 + 1)
            media_plan_clean = media_plan_clean.map(lambda x: x.strip() if isinstance(x, str) else x)

            try:
                date_text = str(media_plan_clean['Flight duration'][0])
                parsed_ranges = parse_date_range(date_text)
                
                if parsed_ranges:
                    latest_range = max(parsed_ranges, key=lambda r: r[1])
                    # Format for display - you can adjust format as needed
                    date_range = f"{latest_range[0].strftime('%d/%m/%Y')} - {latest_range[1].strftime('%d/%m/%Y')}"
                else:
                    print(f"No valid date ranges found in: '{date_text}'")
                    # Fallback to original text or current date range
                    date_range = date_text
            except Exception as e:
                print(f"Error processing date range: {e}")
                # Fallback to current date
                today = datetime.now()
                date_range = f"{today.strftime('%d/%m/%Y')} - {today.strftime('%d/%m/%Y')}"
            
            # Group calculations
            estimated_impression = media_plan_clean.groupby('Platform')['Estimated Impressions'].sum()
            estimated_clicks = media_plan_clean.groupby('Platform')['Estimated link clicks'].sum()
            estimated_reach = media_plan_clean.groupby('Platform')['Estimated Reach'].sum()
            estimated_gross_spend = media_plan_clean.groupby('Platform')['Gross spend by channel / platform'].sum()
            est_imp_sum = estimated_impression.sum().astype(int)
            est_click_sum = estimated_clicks.sum().astype(int)
            est_reach_sum = float(estimated_reach.sum())
            est_gross_sum = float(estimated_gross_spend.sum())
            
            # Meta calculations (always present)
            est_ctr_meta = media_plan_clean[media_plan_clean['Platform'] == 'Meta (Facebook & Instagram)']['Estimated CTR'].iloc[0] * 100
            est_cpm_meta = media_plan_clean[media_plan_clean['Platform'] == 'Meta (Facebook & Instagram)']['Net CPM'].iloc[0]
            est_imp_meta = int(estimated_impression['Meta (Facebook & Instagram)'])
            est_clicks_meta = int(estimated_clicks['Meta (Facebook & Instagram)'])
            est_reach_meta = int(estimated_reach['Meta (Facebook & Instagram)'])
            est_gross_meta = int(estimated_gross_spend['Meta (Facebook & Instagram)'])
            
            # Check if Pinterest exists in the Platform column
            has_pinterest = 'Pinterest' in media_plan_clean['Platform'].values
            
            if has_pinterest:
                # Pinterest calculations if available
                est_ctr_pin = media_plan_clean[media_plan_clean['Platform'] == 'Pinterest']['Estimated CTR'].iloc[0] * 100
                est_cpm_pin = media_plan_clean[media_plan_clean['Platform'] == 'Pinterest']['Net CPM'].iloc[0]
                est_imp_pin = int(estimated_impression['Pinterest'])
                est_clicks_pin = int(estimated_clicks['Pinterest'])
                est_reach_pin = int(estimated_reach['Pinterest'])
                est_gross_pin = int(estimated_gross_spend['Pinterest'])
                
                # Calculate average values using both platforms
                est_ctr_sum = (est_ctr_meta + est_ctr_pin) / 2
                est_cpm_sum = (est_cpm_meta + est_cpm_pin) / 2
            else:
                # Set Pinterest values to 0 if not present
                est_ctr_pin = est_cpm_pin = 0
                est_imp_pin = est_clicks_pin = est_reach_pin = est_gross_pin = 0
                
                # Use only Meta values for summary calculations
                est_ctr_sum = est_ctr_meta
                est_cpm_sum = est_cpm_meta

            # Performance vs estimate
            perc_imp = ((total_impressions - est_imp_sum) / est_imp_sum) * 100 if est_imp_sum != 0 else 0
            perc_clicks = ((total_clicks - est_click_sum) / est_click_sum) * 100 if est_click_sum != 0 else 0

            # Slide-level data for JSON

        else:
            #estimated_impression = estimated_clicks = estimated_reach = estimated_gross_spend = 0
            date_range = Date_
            est_imp_sum = est_click_sum = est_reach_sum = est_gross_sum = 0
            est_ctr_sum = est_cpm_sum = perc_imp = perc_clicks = 0
            est_imp_meta = est_clicks_meta = est_reach_meta = est_gross_meta = 0
            est_imp_pin = est_clicks_pin = est_reach_pin = est_gross_pin = 0

        data = {
                "Metric": ["Gross Spend", "Impressions", "Clicks", "CTR", "Net CPM"],
                "Actual": [total_gross_spend, total_impressions, total_clicks, total_CTR, total_CPM],
                "Estimated": [est_gross_sum, est_imp_sum, est_click_sum, est_ctr_sum, est_cpm_sum]
            }
        df_comparison = pd.DataFrame(data)
        actual_and_estimate_values = df_comparison.to_json(orient='records', indent=4)

        #=============================================================Prompt filling communication with OpenAI ==================================================================================
        prompt = prompt.replace('<<meta_input>>', meta_input)
        prompt = prompt.replace('<<pin_input>>', pin_input)
        prompt = prompt.replace('<<estimate_and_actual>>', actual_and_estimate_values)
        prompt = prompt.replace('<<chaThank younnel_comparison>>', channel_kpi_comparison_json)
        prompt = prompt.replace('<<online_and_instore>>', online_offline_json)
        prompt = prompt.replace('<<audience_data>>', audience_json)
        campaign_commentary = self.gpt4_completion(prompt)                                                                              #Run gpt function and save the output
        commentary_json = json.loads(campaign_commentary)
        print("Commentary JSON:", commentary_json)  # Debugging output

        #=============================================================Call graph plotting function  ==================================================================================

        self.plot_figures(est_imp_sum, total_impressions, "impressions_chart.png")
        self.plot_figures(est_click_sum, total_clicks, "clicks_chart.png")

        #=============================================================replacement dictionary==================================================================================
        self.replacements = {
            "{Brand Name}": Brand_Name,
            "{Campaign Name}": Campaign_Name,
            "{Date_}": Date_,
            "{Campaign Dates}": Date_,
            "{Flight_dates}": date_range,
            "{Channels}": channel_list_str,  # Dynamically formatted channel list
            "{Gross_Spend}": f"{int(total_gross_spend):,}",  # Formatted with commas and two decimal places
            "{Impressions}": f"{total_impressions:,}",  # Formatted with commas
            "{Clicks}": f"{total_clicks:,}",
            "{Net_CPM}": f"{total_CPM:,.2f}",  # Cost per 1000 impressions
            "{CTR}": f"{total_CTR:.2f}",  # Click-through rate percentage
            "{gross_est}": f"{int(est_gross_sum):,}",
            "{imp_est}": f"{est_imp_sum:,}",
            "{reach_est}": f"{int(est_reach_sum):,}",
            "{click_est}": f"{est_click_sum:,}",
            "{net_cpm_est}": f"{est_cpm_sum:,.2f}",
            "{ctr_est}": f"{est_ctr_sum:,.2f}",
            "{perc_imp}": f"{perc_imp:+,.1f}",
            "{perc_clicks}": f"{perc_clicks:+,.1f}",
            "{Brand Revenue}": f"{total_brand_revenue:,.2f}",
            "{Brand ROAS}": f"{total_ROAS:.2f}",
            "{Brand ROI}": f"{total_ROI:.2f}",
            "{FSKU Revenue}": f"{total_FSKU_rev:,.2f}",
            "{FSKU ROAS}": f"{total_FSKU_ROAS:.2f}",
            "{FSKU ROI}": f"{total_FSKU_ROI:.2f}",
            "{overall engagement performance}": commentary_json['overall_engagement_performance'],
            "{overall sales, revenue and ROI performance}": commentary_json['overall_sales_revenue_ROI_performance'],
            "{overall reach performance}": commentary_json['overall_reach_performance'],
            "{unique and estimated values}": commentary_json['estimated_and_actual_values_comparision'],
            "{channel_commentary}": commentary_json['channel_kpi_comparison'],
            "{online and instore comparison}": commentary_json['instore_and_online_comparison'],
            "{overall audience performance}": commentary_json['overall_audience_performance'],
            "{current brand shoppers' performance​}": commentary_json['current_brand_shoppers_performance'],
            "{acquisitional shoppers' performance​}": commentary_json['acquisitional_shoppers_performance'],
            "{online purchase behaviour​}": commentary_json['online_purchase_behavior'],
            "{overall campaign performance}": commentary_json['overall_campaign_performance'],
            "{Meta_Impressions}": f"{int(impressions_meta):,}",
            "{meta_ctr}": f"{CTR_meta:,.2f}",
            "{Meta_Net_CPM}": f"{net_CPM_meta:,.2f}",
            "{Meta_clicks}": f"{int(clicks_meta):,}",
            "{Meta_Gross_Spend}": f"{int(gross_spend_meta):,}",
            "{Meta_Revenue}": f"{brand_revenue_meta:,.2f}",
            "{Meta_Brand_ROAS}": f"{brand_ROAS_meta:,.2f}",
            "{Meta_Brand_ROI}": f"{brand_ROI_meta:,.2f}",
            "{Meta_FSKU_Revenue}": f"{total_FSKU_rev:,.2f}",
            "{Meta_FSKU_ROAS}": f"{total_FSKU_ROAS:,.2f}",
            "{Meta_FSKU_ROI}": f"{total_FSKU_ROI:,.2f}",
            "{Pin_Impressions}": f"{impressions_pin:,}",
            "{Pin_CTR}": f"{CTR_pin:,.2f}",
            "{Pin_Net_CPM}": f"{net_CPM_pin:,.2f}",
            "{Pin_Clicks}": f"{clicks_pin:,}",
            "{pin_gross_spend}": f"{int(gross_spend_pin):,}",
            "{pin_brand_revenue}": f"{brand_revenue_pin:,.2f}",
            "{pin_brand_ROS}": f"{brand_ROAS_pin:,.2f}",
            "{pin_brand_ROI}": f"{brand_ROI_pin:,.2f}",
            "{Online_rev}": f"{total_brand_online_rev:,.2f}",
            "{Instore_rev}": f"{total_brand_instore_rev:,.2f}",
            "{meta_est_spend}": f"{est_gross_meta:,}",
            "{meta_est_imp}": f"{est_imp_meta:,}",
            "{meta_est_reach}": f"{est_reach_meta:,}",
            "{meta_est_clicks}": f"{est_clicks_meta:,}",
            "{meta_est_CTR}": f"{est_ctr_meta:,}%" if media_plan is not None else "",
            "{meta_est_CPM}": f"{est_cpm_meta:,}" if media_plan is not None else "",
            "{Meta_brand_online_revenue}": f"{brand_online_revenue_meta:,}",
            "{Meta_brand_instore_revenue}": f"{brand_instore_revenue_meta:,}",
            "{FSKU_online_revenue_meta}": f"{FSKU_online_revenue_meta:,}",
            "{FSKU_instore_revenue_meta}":f"{FSKU_instore_revenue_meta:,}",
            "{pin_est_spend}": f"{est_gross_pin:,}",
            "{pin_est_imp}": f"{est_imp_pin:,}",
            "{pin_est_reach}": f"{est_reach_pin:,}",
            "{pin_est_clicks}": f"{est_clicks_pin:,}",
            "{pin_est_CTR}": f"{est_ctr_pin:,}"if media_plan is not None else "",
            "{pin_est_CPM}": f"{est_cpm_pin:,}"if media_plan is not None else "",
            "{pin_brand_online_revenue}": f"{brand_online_revenue_pin:,.2f}",
            "{pin_brand_instore_revenue}": f"{brand_instore_revenue_pin:,.2f}",

        }

        for i, row in merged_audience_metrics.iterrows():
            idx = i + 1  # To match placeholder numbering starting at 1
            self.replacements.update({
                f"{{aud_{idx}}}": row['Ad_Set_Name'],
                f"{{Net_spend_{idx}}}": f"{row['Net_Spend']:.2f}",
                f"{{Imp_{idx}}}": f"{int(row['Impressions']):,}",
                f"{{Reach_{idx}}}": "",
                f"{{Freq_{idx}}}": "",  # Add logic if you have frequency data
                f"{{Clicks_{idx}}}": f"{int(row['Clicks']):,}",
                f"{{CTR_{idx}}}": f"{row['CTR']:.2f}%",
                f"{{NetCPM_{idx}}}": f"{row['Net_CPM']:.2f}",
                f"{{Revenue_{idx}}}": f"{row['Brand_Revenue']:.2f}",
                f"{{ROAS_{idx}}}": f"{row['Brand_ROI']:.2f}"

            })

        for i, row in combined_meta_df.iterrows():
            idx = i + 1  # To match placeholder numbering starting at 1
            self.replacements.update({
                f"{{m_aud_{idx}}}": row['Ad_Set_Name'],
                f"{{m_Net_spend_{idx}}}": f"{row['Net_Spend']:.2f}",
                f"{{m_Imp_{idx}}}": f"{int(row['Impressions']):,}",
                f"{{m_Freq_{idx}}}": "",  # Add logic if you have frequency data
                f"{{m_Reach_{idx}}}": "",  # Add logic if you have frequency data
                f"{{m_Clicks_{idx}}}": f"{int(row['Clicks']):,}",
                f"{{m_CTR_{idx}}}": f"{row['CTR']:.4f}%",
                f"{{m_NetCPM_{idx}}}": f"{row['Net_CPM']:.2f}",
                f"{{m_Revenue_{idx}}}": f"{row['Brand_Revenue']:.2f}",
                f"{{m_ROAS_{idx}}}": f"{row['Brand_ROI']:.4f}"

            })

        for i, row in combined_pin_df.iterrows():
            idx = i + 1  # To match placeholder numbering starting at 1
            self.replacements.update({
                f"{{p_aud_{idx}}}": row['Ad_Set_Name'],
                f"{{p_Net_spend_{idx}}}": f"{row['Net_Spend']:.2f}",
                f"{{p_Imp_{idx}}}": f"{int(row['Impressions']):,}",
                f"{{p_Reach_{idx}}}": f"{int(row['Reach']):,}",
                f"{{p_Freq_{idx}}}": "",  # Add logic if you have frequency data
                f"{{p_Clicks_{idx}}}": f"{int(row['Clicks']):,}",
                f"{{p_CTR_{idx}}}": f"{row['CTR']:.2f}%",
                f"{{p_NetCPM_{idx}}}": f"{row['Net_CPM']:.2f}",
                f"{{p_Revenue_{idx}}}": f"{row['Brand_Revenue']:.2f}",
                f"{{p_ROI_{idx}}}": f"{row['Brand_ROI']:.2f}"

            })
        return self.replacements, self.is_pin

st.markdown("""
    <style>
    /* Set background color for the entire app */
    .stApp {
        background-color: #000;
        color: #fff;
    }

    /* Style markdown text */
    .block-container {
        color: #fff;
    }

    /* Style file uploader label */
    label {
        background-color: #e1ad01;
        color: #000 !important;
    }

    /* Style tab buttons (if using st.tabs) */
    div[data-baseweb="tab"] button {
        background-color: #e1ad01 !important;
        color: black !important;
        font-weight: bold;
    }

    /* Style buttons */
    .stButton>button {
        background-color: #e1ad01;
        color: black;
        font-weight: bold;
        border: none;
    }

    /* Style file upload input text */
    .stFileUploader label span {
        color: white;
    }
    </style>
""", unsafe_allow_html=True)


def load_logo_base64(file_path):
    with open(file_path, "rb") as f:
        return base64.b64encode(f.read()).decode()

logo_base64 = load_logo_base64("smg2.jpeg")
st.markdown(
    f"""
    <div style="text-align: center;">
        <img src="data:image/png;base64,{logo_base64}" width="100">
    </div>
    """,
    unsafe_allow_html=True,
)



# Initialize session state variables if they don't exist
if "reset_state" not in st.session_state:
    st.session_state.reset_state = False

if "meta_file" not in st.session_state:
    st.session_state.meta_file = None

if "pinterest_file" not in st.session_state:
    st.session_state.pinterest_file = None

if "media_plan_file" not in st.session_state:
    st.session_state.media_plan_file = None

st.title("📊 Post-Campaign Report Generator")

st.markdown("""
Upload your campaign files below. At minimum, a Meta Excel/csv file is required. Ensure data columns align to column name requirements
""")

# Handle file uploads with unique keys that change when reset is pressed
upload_key_suffix = f"_{hash(st.session_state.reset_state)}"

meta_file = st.file_uploader("Meta Excel File (Required)", 
                            type=["xlsx"], 
                            key=f"meta_uploader{upload_key_suffix}")

pinterest_file = st.file_uploader("Pinterest CSV File (Optional)", 
                                 type=["csv"], 
                                 key=f"pinterest_uploader{upload_key_suffix}")

media_plan_file = st.file_uploader("Media Plan Excel File (Optional)", 
                                  type=["xlsx"], 
                                  key=f"media_plan_uploader{upload_key_suffix}")

# Store uploaded files in session state
if meta_file is not None:
    st.session_state.meta_file = meta_file

if pinterest_file is not None:
    st.session_state.pinterest_file = pinterest_file

if media_plan_file is not None:
    st.session_state.media_plan_file = media_plan_file

col1, col2, col3, col4 = st.columns(4)

with col1:
    if st.button("Generate PowerPoint Report"):
        if not st.session_state.meta_file:
            st.error("Please upload at least a Meta file.")
        else:
            with tempfile.TemporaryDirectory() as tmpdir:
                # Save uploaded files temporarily
                meta_path = os.path.join(tmpdir, st.session_state.meta_file.name)
                # Reset file position to beginning before reading
                st.session_state.meta_file.seek(0)
                with open(meta_path, "wb") as f:
                    f.write(st.session_state.meta_file.read())

                pin_path = None
                if st.session_state.pinterest_file:
                    pin_path = os.path.join(tmpdir, st.session_state.pinterest_file.name)
                    st.session_state.pinterest_file.seek(0)
                    with open(pin_path, "wb") as f:
                        f.write(st.session_state.pinterest_file.read())

                media_path = None
                if st.session_state.media_plan_file:
                    media_path = os.path.join(tmpdir, st.session_state.media_plan_file.name)
                    st.session_state.media_plan_file.seek(0)
                    with open(media_path, "wb") as f:
                        f.write(st.session_state.media_plan_file.read())

                # Provide internal file paths (ensure these files exist in your project directory)
                prompt = "prompt.txt"
                ppt_template = "automation_template_v3.pptx"

                prompt_path = os.path.join(tmpdir, "prompt.txt")
                with open(prompt, "rb") as src, open(prompt_path, "wb") as dst:
                    dst.write(src.read())

                ppt_template_path = os.path.join(tmpdir, "template.pptx")
                with open(ppt_template, "rb") as src, open(ppt_template_path, "wb") as dst:
                    dst.write(src.read())

                # Run processing
                with st.spinner("⏳ Generating your PowerPoint report..."):
                    extractor = DataExtractor(meta_path, pin_path, media_path, prompt_path)
                    replacements, is_pin = extractor.extract_values()

                    output_path = os.path.join(tmpdir, "automated_presentation.pptx")
                    ppt = PowerPointProcessor(ppt_template_path)
                    ppt.replace_placeholders(replacements, is_pin, output_path)

                # Download output
                with open(output_path, "rb") as file:
                    st.success("✅ PowerPoint report generated!")
                    st.download_button(
                        label="📥 Download Report",
                        data=file,
                        file_name="automated_presentation.pptx",
                        mime="application/vnd.openxmlformats-officedocument.presentationml.presentation"
                    )
with col4:
    if st.button("Reset Uploads"):
        # Toggle the reset state to force file uploader widgets to create new instances
        st.session_state.reset_state = not st.session_state.reset_state
        # Clear stored files
        st.session_state.meta_file = None
        st.session_state.pinterest_file = None
        st.session_state.media_plan_file = None
        st.rerun()
